"""
PPT Generator with Claude AI
웹 기반 PPT 생성 도구 - Claude API를 사용하여 스마트한 프레젠테이션 생성
"""

import streamlit as st
from pptx import Presentation
from pptx.util import Pt
import anthropic
import json
import io
import re
import os
from datetime import datetime
from dotenv import load_dotenv

# 환경 변수 로드
load_dotenv()

# 페이지 설정
st.set_page_config(
    page_title="AI PPT Generator - Powered by Kwak Dooil",
    page_icon="🎨",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# Modern CSS 스타일
st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap');
    
    /* 전체 폰트 및 배경 */
    * {
        font-family: 'Inter', -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif;
    }
    
    .stApp {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        background-attachment: fixed;
    }
    
    .main {
        background: rgba(255, 255, 255, 0.95);
        border-radius: 20px;
        padding: 2rem;
        margin: 1rem;
        backdrop-filter: blur(10px);
        box-shadow: 0 25px 50px -12px rgba(0, 0, 0, 0.25);
    }
    
    /* 헤더 스타일 */
    h1, .stMarkdown h1, [data-testid="stMarkdown"] h1 {
        color: #667eea !important;
        font-weight: 700 !important;
        font-size: 3rem !important;
        text-align: center !important;
        margin-bottom: 0.5rem !important;
        animation: fadeInDown 0.8s ease;
        text-shadow: 2px 2px 4px rgba(0,0,0,0.1);
        display: block !important;
        visibility: visible !important;
    }
    
    h2 {
        color: #1a202c;
        font-weight: 600;
        border-bottom: 3px solid #667eea;
        padding-bottom: 0.5rem;
        margin-top: 2rem;
    }
    
    h3 {
        color: #2d3748;
        font-weight: 500;
    }
    
    /* 서브타이틀 */
    .subtitle {
        text-align: center;
        color: #718096;
        font-size: 1.1rem;
        margin-bottom: 2rem;
        animation: fadeInUp 0.8s ease;
    }
    
    /* 카드 스타일 */
    .css-1r6slb0, .css-12w0qpk {
        background: white;
        border-radius: 15px;
        padding: 1.5rem;
        box-shadow: 0 10px 30px rgba(0, 0, 0, 0.1);
        border: 1px solid rgba(102, 126, 234, 0.1);
        transition: all 0.3s ease;
    }
    
    .css-1r6slb0:hover, .css-12w0qpk:hover {
        transform: translateY(-5px);
        box-shadow: 0 15px 40px rgba(102, 126, 234, 0.2);
    }
    
    /* 버튼 스타일 */
    .stButton > button {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        color: white;
        border: none;
        padding: 0.75rem 2rem;
        font-weight: 600;
        border-radius: 50px;
        transition: all 0.3s ease;
        box-shadow: 0 4px 15px rgba(102, 126, 234, 0.4);
        width: 100%;
        font-size: 1.1rem;
    }
    
    .stButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 6px 20px rgba(102, 126, 234, 0.5);
    }
    
    .stButton > button:active {
        transform: translateY(0);
    }
    
    /* Primary 버튼 특별 스타일 */
    div[data-testid="stButton"] button[kind="primary"] {
        background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%);
        font-size: 1.2rem;
        padding: 1rem 2.5rem;
    }
    
    /* 다운로드 버튼 */
    .stDownloadButton > button {
        background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%);
        color: white;
        border: none;
        padding: 0.75rem 2rem;
        font-weight: 600;
        border-radius: 50px;
        transition: all 0.3s ease;
        box-shadow: 0 4px 15px rgba(79, 172, 254, 0.4);
    }
    
    .stDownloadButton > button:hover {
        transform: translateY(-2px);
        box-shadow: 0 6px 20px rgba(79, 172, 254, 0.5);
    }
    
    /* 입력 필드 */
    .stTextInput > div > div > input,
    .stTextArea > div > div > textarea,
    .stSelectbox > div > div > div {
        border: 2px solid #e2e8f0;
        border-radius: 10px;
        padding: 0.75rem;
        transition: all 0.3s ease;
        background: white;
    }
    
    .stTextInput > div > div > input:focus,
    .stTextArea > div > div > textarea:focus {
        border-color: #667eea;
        box-shadow: 0 0 0 3px rgba(102, 126, 234, 0.1);
        outline: none;
    }
    
    /* 슬라이더 */
    .stSlider > div > div > div {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
    }
    
    /* 탭 스타일 */
    .stTabs [data-baseweb="tab-list"] {
        background: rgba(102, 126, 234, 0.05);
        border-radius: 10px;
        padding: 0.25rem;
        gap: 0.5rem;
    }
    
    .stTabs [data-baseweb="tab"] {
        border-radius: 8px;
        padding: 0.5rem 1rem;
        background: transparent;
        color: #4a5568;
        transition: all 0.3s ease;
    }
    
    .stTabs [aria-selected="true"] {
        background: white;
        color: #667eea;
        box-shadow: 0 2px 10px rgba(0, 0, 0, 0.1);
    }
    
    /* 사이드바 */
    .css-1d391kg {
        background: linear-gradient(180deg, #667eea 0%, #764ba2 100%);
    }
    
    .css-1d391kg .stButton > button {
        background: rgba(255, 255, 255, 0.2);
        backdrop-filter: blur(10px);
        border: 1px solid rgba(255, 255, 255, 0.3);
    }
    
    .css-1d391kg .stButton > button:hover {
        background: rgba(255, 255, 255, 0.3);
    }
    
    /* Success/Warning/Error 메시지 */
    .stSuccess, .stWarning, .stError, .stInfo {
        border-radius: 10px;
        padding: 1rem;
        border-left: 4px solid;
        animation: slideInRight 0.5s ease;
    }
    
    /* 분할선 */
    hr {
        border: none;
        height: 2px;
        background: linear-gradient(90deg, transparent, #667eea, transparent);
        margin: 2rem 0;
    }
    
    /* 푸터 */
    .footer {
        text-align: center;
        padding: 2rem;
        margin-top: 3rem;
        background: rgba(102, 126, 234, 0.05);
        border-radius: 15px;
        color: #4a5568;
    }
    
    /* 애니메이션 */
    @keyframes fadeInDown {
        from {
            opacity: 0;
            transform: translateY(-20px);
        }
        to {
            opacity: 1;
            transform: translateY(0);
        }
    }
    
    @keyframes fadeInUp {
        from {
            opacity: 0;
            transform: translateY(20px);
        }
        to {
            opacity: 1;
            transform: translateY(0);
        }
    }
    
    @keyframes slideInRight {
        from {
            opacity: 0;
            transform: translateX(-20px);
        }
        to {
            opacity: 1;
            transform: translateX(0);
        }
    }
    
    /* 숫자 입력 */
    .stNumberInput > div > div > input {
        border: 2px solid #e2e8f0;
        border-radius: 10px;
        padding: 0.75rem;
    }
    
    /* Expander */
    .streamlit-expanderHeader {
        background: rgba(102, 126, 234, 0.05);
        border-radius: 10px;
        border: 1px solid rgba(102, 126, 234, 0.2);
    }
    
    /* 메트릭 카드 */
    [data-testid="metric-container"] {
        background: linear-gradient(135deg, rgba(102, 126, 234, 0.1) 0%, rgba(118, 75, 162, 0.1) 100%);
        border: 1px solid rgba(102, 126, 234, 0.2);
        padding: 1rem;
        border-radius: 10px;
        box-shadow: 0 2px 10px rgba(0, 0, 0, 0.05);
    }
    
    /* 스피너 */
    .stSpinner > div {
        border-top-color: #667eea !important;
    }
    
    /* 코드 블록 */
    .stCodeBlock {
        border-radius: 10px;
        background: #1a202c;
    }
    </style>
    """, unsafe_allow_html=True)

# 초기화
if 'api_key' not in st.session_state:
    st.session_state.api_key = os.getenv('CLAUDE_API_KEY', '')

if 'generated_slides' not in st.session_state:
    st.session_state.generated_slides = []

def analyze_topic_and_suggest_template(topic):
    """주제 분석하여 최적 템플릿과 구조 제안"""
    topic_lower = topic.lower()
    
    # 키워드 기반 자동 템플릿 선택
    if any(word in topic_lower for word in ['마케팅', '판매', '브랜딩', '광고', '홍보']):
        return {'template': 'modern', 'color': 'orange-red', 'type': '마케팅'}
    elif any(word in topic_lower for word in ['기술', '개발', 'ai', '인공지능', 'it', '프로그래밍']):
        return {'template': 'dark', 'color': 'blue-purple', 'type': '기술'}
    elif any(word in topic_lower for word in ['교육', '학습', '강의', '연수', '교수법']):
        return {'template': 'minimal', 'color': 'green-blue', 'type': '교육'}
    elif any(word in topic_lower for word in ['회사', '기업', '비즈니스', '전략', '경영']):
        return {'template': 'corporate', 'color': 'blue-purple', 'type': '비즈니스'}
    else:
        return {'template': 'modern', 'color': 'blue-purple', 'type': '일반'}

def generate_slides_with_claude(topic, num_slides, style):
    """Claude API를 사용하여 슬라이드 콘텐츠 생성"""
    client = anthropic.Anthropic(api_key=st.session_state.api_key)
    
    # 주제 분석 및 자동 템플릿 선택
    template_info = analyze_topic_and_suggest_template(topic)
    
    # 자동 선택된 템플릿 적용
    st.session_state['design_template'] = template_info['template']
    st.session_state['color_scheme'] = template_info['color']
    
    prompt = f"""
    당신은 전문 프레젠테이션 디자이너입니다. 다음 정보로 완벽한 프레젠테이션을 만들어주세요:

    📋 프로젝트 정보:
    - 주제: {topic}
    - 슬라이드 수: {num_slides}개
    - 스타일: {style}
    - 추천 템플릿: {template_info['template']} ({template_info['type']} 특화)

    📐 프레젠테이션 구조 가이드:
    - 1번 슬라이드: 타이틀 + 매력적인 부제목
    - 2번 슬라이드: 목차/개요 (전체 흐름 제시)
    - 중간 슬라이드들: 핵심 내용 (논리적 순서)
    - 마지막 슬라이드: 결론/행동 촉구

    🎯 내용 작성 원칙:
    1. 각 슬라이드는 하나의 핵심 메시지에 집중
    2. 제목은 임팩트 있고 명확하게 (10자 이내 권장)
    3. 내용은 3-5개 핵심 포인트로 구성
    4. 전문적이면서도 이해하기 쉽게
    5. 실행 가능한 구체적 내용 포함

    ⚡ 필수 규칙:
    - bullet point 기호(•, -, *, ○ 등) 절대 사용 금지
    - 각 포인트는 줄바꿈(\\n)으로만 구분
    - {style} 스타일에 맞는 톤 유지
    - 숫자, 통계, 구체적 예시 적극 활용

    📊 JSON 형식으로 반환:
    [
        {{
            "title": "매력적인 슬라이드 제목",
            "content": "첫 번째 핵심 내용\\n두 번째 핵심 내용\\n세 번째 핵심 내용"
        }},
        ...
    ]

    💡 예시 (마케팅 주제의 경우):
    {{
        "title": "2025 디지털 마케팅 전략",
        "content": "소셜미디어 중심의 브랜드 스토리텔링\\n개인화된 고객 경험 설계\\nAI 기반 마케팅 자동화 도입\\n성과 측정 및 최적화 시스템 구축"
    }}
    
    주제 '{topic}'에 대해 {num_slides}개의 완벽한 슬라이드를 생성해주세요.
    """
    
    try:
        message = client.messages.create(
            model="claude-3-5-sonnet-20241022",
            max_tokens=4000,
            temperature=0.7,
            system="당신은 전문적인 프레젠테이션 콘텐츠 작성자입니다. 주어진 주제에 대해 체계적이고 논리적인 PPT 슬라이드를 생성합니다.",
            messages=[
                {"role": "user", "content": prompt}
            ]
        )
        
        # 응답에서 JSON 추출
        response_text = message.content[0].text
        
        # JSON 부분만 추출 (```json 태그 제거)
        json_match = re.search(r'\[.*\]', response_text, re.DOTALL)
        if json_match:
            json_str = json_match.group()
            slides_data = json.loads(json_str)
            return slides_data
        else:
            # JSON 형식이 없으면 텍스트를 파싱
            return parse_text_to_slides(response_text, num_slides)
            
    except Exception as e:
        st.error(f"Claude API 오류: {str(e)}")
        return None

def parse_text_to_slides(text, num_slides):
    """텍스트를 슬라이드 형식으로 파싱"""
    lines = text.strip().split('\n')
    slides = []
    current_slide = None
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        # 숫자로 시작하거나 "슬라이드"가 포함된 라인을 제목으로 간주
        if re.match(r'^\d+\.', line) or '슬라이드' in line or len(slides) < num_slides and not current_slide:
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                "title": line.lstrip('0123456789. '),
                "content": ""
            }
        elif current_slide:
            if current_slide["content"]:
                current_slide["content"] += "\n"
            current_slide["content"] += line
    
    if current_slide:
        slides.append(current_slide)
    
    # 슬라이드 개수 맞추기
    while len(slides) < num_slides:
        slides.append({"title": f"슬라이드 {len(slides) + 1}", "content": "내용을 추가해주세요"})
    
    return slides[:num_slides]

def create_slide_preview(title, content, design_template="modern", color_scheme="blue-purple"):
    """Gamma.app 스타일의 실시간 슬라이드 프리뷰 생성"""
    
    # 색상 스킴 정의
    color_schemes = {
        "blue-purple": {"primary": "#667eea", "secondary": "#764ba2", "bg": "linear-gradient(135deg, #667eea 0%, #764ba2 100%)"},
        "green-blue": {"primary": "#11998e", "secondary": "#38ef7d", "bg": "linear-gradient(135deg, #11998e 0%, #38ef7d 100%)"},
        "orange-red": {"primary": "#ff6b6b", "secondary": "#ffa726", "bg": "linear-gradient(135deg, #ff6b6b 0%, #ffa726 100%)"},
        "purple-pink": {"primary": "#667eea", "secondary": "#f093fb", "bg": "linear-gradient(135deg, #667eea 0%, #f093fb 100%)"}
    }
    
    scheme = color_schemes.get(color_scheme, color_schemes["blue-purple"])
    
    # 내용 처리
    if not content:
        content_html = "<div class='content-item'>내용을 입력해주세요</div>"
    else:
        content_lines = content.split('\n')
        formatted_content = []
        
        for line in content_lines:
            line = line.strip()
            if line:
                # bullet point 제거
                cleaned_line = line
                for bullet in ['•', '-', '*', '◦', '○', '▪', '▫', '‣', '⁃']:
                    if cleaned_line.startswith(bullet):
                        cleaned_line = cleaned_line[1:].strip()
                        break
                
                if cleaned_line:
                    formatted_content.append(f"<div class='content-item'>• {cleaned_line}</div>")
        
        content_html = '\n'.join(formatted_content) if formatted_content else "<div class='content-item'>내용을 입력해주세요</div>"
    
    # 프리뷰 HTML 생성
    preview_html = f"""
    <div style="
        background: {scheme['bg']};
        border-radius: 15px;
        padding: 30px;
        min-height: 300px;
        color: white;
        box-shadow: 0 10px 30px rgba(0,0,0,0.2);
        font-family: 'Inter', sans-serif;
        position: relative;
        overflow: hidden;
    ">
        <!-- 장식 요소 -->
        <div style="
            position: absolute;
            top: -50px;
            right: -50px;
            width: 100px;
            height: 100px;
            background: rgba(255,255,255,0.1);
            border-radius: 50%;
        "></div>
        
        <!-- 제목 -->
        <h2 style="
            color: white;
            font-size: 1.8rem;
            font-weight: 700;
            margin-bottom: 20px;
            line-height: 1.2;
            text-shadow: 0 2px 4px rgba(0,0,0,0.3);
        ">{title}</h2>
        
        <!-- 내용 -->
        <div style="
            color: rgba(255,255,255,0.9);
            font-size: 1rem;
            line-height: 1.6;
        ">
            {content_html}
        </div>
        
        <!-- 하단 장식 -->
        <div style="
            position: absolute;
            bottom: 0;
            left: 0;
            right: 0;
            height: 4px;
            background: rgba(255,255,255,0.3);
        "></div>
    </div>
    
    <style>
    .content-item {{
        margin-bottom: 8px;
        padding-left: 10px;
        animation: fadeIn 0.5s ease;
    }}
    
    @keyframes fadeIn {{
        from {{ opacity: 0; transform: translateY(10px); }}
        to {{ opacity: 1; transform: translateY(0); }}
    }}
    </style>
    """
    
    return preview_html

def create_ppt(slides_data, title_size=24, content_size=16, design_template="modern"):
    """PPT 파일 생성 - 고급 디자인 템플릿 지원"""
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    
    prs = Presentation()
    
    # 디자인 템플릿별 색상 설정 (RGB 값으로 변환)
    design_themes = {
        "modern": {"title_color": RGBColor(102, 126, 234), "content_color": RGBColor(45, 55, 72), "bg_color": RGBColor(255, 255, 255)},
        "dark": {"title_color": RGBColor(79, 172, 254), "content_color": RGBColor(255, 255, 255), "bg_color": RGBColor(26, 26, 26)},
        "minimal": {"title_color": RGBColor(40, 167, 69), "content_color": RGBColor(33, 37, 41), "bg_color": RGBColor(248, 249, 250)},
        "corporate": {"title_color": RGBColor(0, 86, 179), "content_color": RGBColor(51, 51, 51), "bg_color": RGBColor(255, 255, 255)}
    }
    
    theme = design_themes.get(design_template, design_themes["modern"])
    
    for slide_data in slides_data:
        # 빈 레이아웃 사용해서 완전히 커스터마이징
        slide_layout = prs.slide_layouts[6]  # 빈 레이아웃
        slide = prs.slides.add_slide(slide_layout)
        
        # 배경 설정
        try:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = theme["bg_color"]
        except:
            pass
        
        # 제목 텍스트박스 생성 (상단, 더 크고 굵게)
        if slide_data.get("title"):
            title_left = prs.slide_width * 0.05  # 5% 여백
            title_top = prs.slide_height * 0.1   # 10% 위치
            title_width = prs.slide_width * 0.9  # 90% 너비
            title_height = prs.slide_height * 0.2 # 20% 높이
            
            title_textbox = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_frame = title_textbox.text_frame
            title_frame.text = slide_data["title"]
            
            # 제목 스타일링
            title_paragraph = title_frame.paragraphs[0]
            title_paragraph.alignment = PP_ALIGN.CENTER
            title_run = title_paragraph.runs[0]
            title_run.font.name = '맑은 고딕'
            title_run.font.size = Pt(title_size + 8)  # 더 크게
            title_run.font.bold = True
            title_run.font.color.rgb = theme["title_color"]
        
        # 내용 텍스트박스 생성 (중앙, bullet point 처리)
        if slide_data.get("content"):
            content_left = prs.slide_width * 0.1   # 10% 여백
            content_top = prs.slide_height * 0.35  # 35% 위치
            content_width = prs.slide_width * 0.8  # 80% 너비
            content_height = prs.slide_height * 0.5 # 50% 높이
            
            content_textbox = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            content_frame = content_textbox.text_frame
            
            # 기존 bullet point 제거 후 재구성
            content_text = slide_data["content"]
            content_lines = content_text.split('\n')
            clean_lines = []
            
            for line in content_lines:
                line = line.strip()
                if line:
                    # 모든 종류의 bullet point 제거
                    cleaned_line = line
                    for bullet in ['•', '-', '*', '◦', '○', '▪', '▫', '‣', '⁃']:
                        if cleaned_line.startswith(bullet):
                            cleaned_line = cleaned_line[1:].strip()
                            break
                    if cleaned_line:
                        clean_lines.append(cleaned_line)
            
            # 설정에 따라 bullet point 추가
            use_bullets = st.session_state.get('use_bullets', True)
            
            if clean_lines:
                # 첫 번째 줄
                first_line = clean_lines[0]
                if use_bullets and not first_line.startswith('•'):
                    content_frame.text = f"• {first_line}"
                else:
                    content_frame.text = first_line
                
                # 나머지 줄들
                for line in clean_lines[1:]:
                    p = content_frame.add_paragraph()
                    if use_bullets and not line.startswith('•'):
                        p.text = f"• {line}"
                    else:
                        p.text = line
                    p.level = 0
            
            # 내용 스타일링
            for paragraph in content_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = '맑은 고딕'
                    run.font.size = Pt(content_size)
                    run.font.color.rgb = theme["content_color"]
                paragraph.space_after = Pt(12)  # 줄 간격
        
        # 하단 장식 라인 추가 (현대적 느낌)
        try:
            line_left = prs.slide_width * 0.1
            line_top = prs.slide_height * 0.9
            line_width = prs.slide_width * 0.8
            line_height = Pt(4)
            
            line_shape = slide.shapes.add_shape(
                1, line_left, line_top, line_width, line_height  # 1 = rectangle
            )
            line_fill = line_shape.fill
            line_fill.solid()
            line_fill.fore_color.rgb = theme["title_color"]
        except:
            pass
    
    # 메모리에 저장
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    
    return ppt_io

# UI 구성
st.title("✨ AI PPT Generator")
st.markdown('<p class="subtitle">지인의 부탁으로 시작한 스마트 프레젠테이션 생성기</p>', unsafe_allow_html=True)

# 메인 컨테이너
main_container = st.container()

# 사이드바 - 설정
with st.sidebar:
    st.markdown("<h2 style='color: white; text-align: center;'>⚙️ 설정</h2>", unsafe_allow_html=True)
    
    # API 키 설정
    with st.expander("🔑 API 설정", expanded=False):
        if st.session_state.api_key:
            st.success("✅ API Key 로드됨")
        else:
            st.warning("⚠️ .env 파일에 CLAUDE_API_KEY를 설정하세요")
    
        api_key_input = st.text_input(
            "Claude API Key",
            value=st.session_state.api_key,
            type="password",
            help="Claude API 키를 입력하세요 (비어있으면 .env 파일 사용)",
            placeholder="sk-ant-api03-..."
        )
        
        if api_key_input != st.session_state.api_key:
            st.session_state.api_key = api_key_input
    
    st.divider()
    
    # 폰트 크기 설정
    with st.expander("📝 폰트 설정", expanded=True):
        title_font_size = st.slider("제목 폰트 크기", 18, 36, 24, help="슬라이드 제목 크기")
        content_font_size = st.slider("내용 폰트 크기", 12, 24, 16, help="슬라이드 본문 크기")
    
    st.divider()
    
    # 사용 가이드
    with st.expander("📖 사용 가이드", expanded=False):
        st.markdown("""
        ### 🚀 빠른 시작
        1. **주제 입력** - PPT 주제를 입력하세요
        2. **슬라이드 수** - 원하는 개수 선택
        3. **스타일 선택** - 프레젠테이션 스타일
        4. **생성** - 'PPT 생성' 버튼 클릭
        5. **편집** - 내용 확인 및 수정
        6. **다운로드** - PPT 파일 저장
        
        ### 💡 팁
        - 구체적인 주제일수록 좋은 결과
        - 생성 후 자유롭게 편집 가능
        - JSON 백업으로 데이터 보관
        """)

# 메인 컨텐츠
with main_container:
    # 입력 섹션과 미리보기 섹션을 카드 스타일로 구분
    col1, col2 = st.columns([1, 1], gap="large")
    
    with col1:
        st.markdown("<h2 style='margin-top: 0;'>🚀 AI 스마트 PPT 생성기</h2>", unsafe_allow_html=True)
        st.markdown("**간단한 설명만으로 완벽한 프레젠테이션을 만들어 드립니다!**")
    
        # 스마트 입력 카드
        with st.container():
            topic = st.text_area(
                "💬 어떤 프레젠테이션을 만들고 싶으신가요?",
                placeholder="자연스럽게 말하듯이 입력하세요!\n\n예시:\n• '우리 팀 2025년 마케팅 전략을 경영진에게 발표해야 해'\n• 'AI 기술을 고객들에게 쉽게 설명하는 PPT 필요해'\n• '신입사원들을 위한 회사 소개 자료 만들어줘'\n• '온라인 교육용 파이썬 기초 강의 자료 부탁해'",
                height=140,
                help="🎯 AI가 자동으로 최적의 구조, 디자인, 내용을 생성합니다!"
            )
    
            # 간단한 추가 옵션
            col1_1, col1_2, col1_3 = st.columns(3)
            
            with col1_1:
                num_slides = st.selectbox(
                    "📊 분량",
                    [5, 8, 10, 12, 15, 20],
                    index=2,
                    format_func=lambda x: f"{x}장",
                    help="슬라이드 개수"
                )
    
            with col1_2:
                target_audience = st.selectbox(
                    "👥 대상",
                    ["팀원/동료", "고객/클라이언트", "경영진", "학생/교육생", "일반 대중"],
                    help="누구를 위한 프레젠테이션인가요?"
                )
                
            with col1_3:
                presentation_time = st.selectbox(
                    "⏰ 발표 시간",
                    ["5분", "10분", "15분", "20분", "30분", "60분+"],
                    index=2,
                    help="예상 발표 시간"
                )
            
            # 자동 스타일 결정 (숨김 처리 - AI가 주제 분석해서 자동 결정)
            style = "비즈니스 (전문적)"  # 기본값, 실제로는 analyze_topic_and_suggest_template에서 결정
            
            # 고급 디자인 옵션
            with st.expander("🎯 고급 디자인 옵션", expanded=False):
                col_design1, col_design2 = st.columns(2)
                
                with col_design1:
                    st.session_state['design_template'] = st.selectbox(
                        "📐 디자인 템플릿",
                        ["modern", "dark", "minimal", "corporate"],
                        index=0,
                        format_func=lambda x: {
                            "modern": "🔮 모던 (그라데이션)",
                            "dark": "🌙 다크 (어두운 테마)",
                            "minimal": "🍃 미니멀 (깔끔한)",
                            "corporate": "💼 기업용 (전문적)"
                        }[x],
                        help="PPT의 전체적인 디자인 테마를 선택하세요",
                        key="design_template_selector"
                    )
                    
                    st.session_state['use_bullets'] = st.checkbox(
                        "📍 Bullet Point 사용", 
                        value=True,
                        help="체크 해제 시 bullet point 없이 깔끔한 텍스트로 생성",
                        key="bullets_checkbox"
                    )
                
                with col_design2:
                    st.session_state['layout_style'] = st.selectbox(
                        "📄 레이아웃 스타일",
                        ["standard", "centered", "split", "image-text"],
                        index=0,
                        format_func=lambda x: {
                            "standard": "📋 표준 (제목+내용)",
                            "centered": "🎯 중앙 정렬",
                            "split": "⚡ 분할 레이아웃",
                            "image-text": "🖼️ 이미지+텍스트"
                        }[x],
                        help="슬라이드의 레이아웃 구조를 선택하세요",
                        key="layout_selector"
                    )
                    
                    st.session_state['color_scheme'] = st.selectbox(
                        "🎨 색상 조합",
                        ["blue-purple", "green-blue", "orange-red", "purple-pink"],
                        index=0,
                        format_func=lambda x: {
                            "blue-purple": "💙 블루-퍼플",
                            "green-blue": "💚 그린-블루", 
                            "orange-red": "🧡 오렌지-레드",
                            "purple-pink": "💜 퍼플-핑크"
                        }[x],
                        help="프레젠테이션의 주요 색상 테마를 선택하세요",
                        key="color_scheme_selector"
                    )
    
            # AI 정보 표시
            if topic:
                template_info = analyze_topic_and_suggest_template(topic)
                st.info(f"🤖 **AI 분석 결과**: {template_info['type']} 카테고리로 분류되어 **{template_info['template']}** 템플릿과 **{template_info['color']}** 색상을 자동 선택했습니다!")
            
            # 원클릭 마법 생성 버튼
            st.markdown("<br>", unsafe_allow_html=True)
            
            # 더 임팩트 있는 버튼
            button_col1, button_col2, button_col3 = st.columns([0.2, 0.6, 0.2])
            with button_col2:
                if st.button("✨ AI 마법으로 PPT 완성하기! ✨", type="primary", use_container_width=True):
                    if topic:
                        # 더 상세한 로딩 메시지
                        progress_placeholder = st.empty()
                        
                        with st.spinner(""):
                            progress_placeholder.info("🔍 주제 분석 중...")
                            import time
                            time.sleep(0.5)
                            
                            progress_placeholder.info("🎨 최적 디자인 템플릿 선택 중...")
                            time.sleep(0.5)
                            
                            progress_placeholder.info("📝 전문가급 콘텐츠 생성 중...")
                            
                            # 더 자세한 프롬프트로 슬라이드 생성
                            enhanced_prompt = f"""
                            {topic}
                            
                            대상: {target_audience}
                            발표시간: {presentation_time}
                            슬라이드수: {num_slides}장
                            
                            위 정보를 종합하여 완벽한 프레젠테이션을 만들어주세요.
                            """
                            
                            slides = generate_slides_with_claude(enhanced_prompt, num_slides, style)
                            
                            progress_placeholder.info("🎯 최종 검토 및 완성 중...")
                            time.sleep(0.3)
                            
                        progress_placeholder.empty()
                        
                        if slides:
                            st.session_state.generated_slides = slides
                            st.balloons()
                            st.success(f"🎉 **완벽한 {len(slides)}장 PPT가 완성되었습니다!** 오른쪽에서 미리보기와 편집을 해보세요.")
                            
                            # 성공 정보를 간단하게 표시
                            st.info(f"📊 **{len(slides)}장 슬라이드** | 🎨 **{st.session_state.get('design_template', 'modern').title()}** 템플릿 | ⏱️ **{presentation_time}** 발표용")
                                
                        else:
                            st.error("😔 잠시 문제가 발생했습니다. 주제를 더 구체적으로 입력하고 다시 시도해주세요.")
                    else:
                        st.warning("⚠️ 먼저 어떤 프레젠테이션을 만들고 싶은지 알려주세요!")

    with col2:
        st.markdown("<h2 style='margin-top: 0;'>👁️ 미리보기 및 편집</h2>", unsafe_allow_html=True)
    
        if st.session_state.generated_slides:
            # 진행 상태 표시
            progress_text = f"총 {len(st.session_state.generated_slides)}개 슬라이드"
            st.info(f"📑 {progress_text}")
            
            # 슬라이드 네비게이션 (탭 대신 페이지네이션 사용)
            if 'current_slide' not in st.session_state:
                st.session_state.current_slide = 0
            
            total_slides = len(st.session_state.generated_slides)
            
            # 네비게이션 컨트롤
            col_prev, col_info, col_next = st.columns([1, 2, 1])
            
            with col_prev:
                if st.button("⬅️ 이전", disabled=st.session_state.current_slide <= 0):
                    st.session_state.current_slide -= 1
                    st.rerun()
            
            with col_info:
                # 슬라이드 선택 selectbox
                st.session_state.current_slide = st.selectbox(
                    "슬라이드 선택",
                    range(total_slides),
                    index=st.session_state.current_slide,
                    format_func=lambda x: f"📄 슬라이드 {x+1} / {total_slides}",
                    key="slide_selector"
                )
            
            with col_next:
                if st.button("➡️ 다음", disabled=st.session_state.current_slide >= total_slides - 1):
                    st.session_state.current_slide += 1
                    st.rerun()
            
            # 현재 슬라이드 편집 및 실시간 프리뷰
            current_idx = st.session_state.current_slide
            slide = st.session_state.generated_slides[current_idx]
            
            st.markdown(f"### 📝 슬라이드 {current_idx + 1} 편집")
            
            # 편집과 프리뷰를 나란히 배치
            edit_col, preview_col = st.columns([1, 1])
            
            with edit_col:
                st.markdown("**✏️ 편집**")
                # 편집 가능한 필드
                edited_title = st.text_input(
                    "제목",
                    value=slide["title"],
                    key=f"title_{current_idx}_edit",
                    help="슬라이드 제목을 입력하세요"
                )
                
                edited_content = st.text_area(
                    "내용",
                    value=slide["content"],
                    height=250,
                    key=f"content_{current_idx}_edit",
                    help="슬라이드 내용을 입력하세요"
                )
                
                # 수정사항 저장
                st.session_state.generated_slides[current_idx]["title"] = edited_title
                st.session_state.generated_slides[current_idx]["content"] = edited_content
            
            with preview_col:
                st.markdown("**👁️ 실시간 프리뷰**")
                # 간단한 텍스트 기반 프리뷰 (HTML 문제 방지)
                
                # 색상 스킴에 따른 배경색 선택
                color_schemes = {
                    "blue-purple": "#667eea",
                    "green-blue": "#11998e", 
                    "orange-red": "#ff6b6b",
                    "purple-pink": "#667eea"
                }
                
                bg_color = color_schemes.get(st.session_state.get('color_scheme', 'blue-purple'), "#667eea")
                
                # 내용 정리
                display_title = edited_title or "제목을 입력하세요"
                display_content = edited_content or "내용을 입력하세요"
                
                # bullet point 제거 및 정리
                if display_content and display_content != "내용을 입력하세요":
                    content_lines = []
                    for line in display_content.split('\n'):
                        line = line.strip()
                        if line:
                            # bullet point 제거
                            for bullet in ['•', '-', '*', '◦', '○']:
                                if line.startswith(bullet):
                                    line = line[1:].strip()
                                    break
                            if line:
                                content_lines.append(f"• {line}")
                    formatted_content = '\n'.join(content_lines)
                else:
                    formatted_content = "• 내용을 입력하세요"
                
                # 프리뷰 카드
                st.markdown(f"""
                <div style="
                    background: linear-gradient(135deg, {bg_color} 0%, #764ba2 100%);
                    border-radius: 15px;
                    padding: 25px;
                    min-height: 280px;
                    color: white;
                    box-shadow: 0 10px 25px rgba(0,0,0,0.15);
                    font-family: 'Inter', sans-serif;
                ">
                    <h3 style="
                        color: white;
                        font-size: 1.4rem;
                        font-weight: 700;
                        margin-bottom: 15px;
                        text-shadow: 0 2px 4px rgba(0,0,0,0.3);
                    ">{display_title}</h3>
                    
                    <div style="
                        color: rgba(255,255,255,0.9);
                        font-size: 0.95rem;
                        line-height: 1.5;
                        white-space: pre-line;
                    ">{formatted_content}</div>
                </div>
                """, unsafe_allow_html=True)
            
            # 다운로드 섹션
            st.markdown("<hr style='margin: 2rem 0;'>", unsafe_allow_html=True)
            st.markdown("<h3>💾 다운로드 옵션</h3>", unsafe_allow_html=True)
            
            col_a, col_b = st.columns(2, gap="medium")
        
            with col_a:
                # 원클릭 PPT 완성 다운로드
                if st.button("🚀 완벽한 PPT 바로 다운로드", type="primary", use_container_width=True):
                    with st.spinner("🎨 디자인 적용하고 파일 생성 중..."):
                        ppt_file = create_ppt(
                            st.session_state.generated_slides,
                            title_font_size,
                            content_font_size,
                            st.session_state.get('design_template', 'modern')
                        )
                        
                        # 더 의미있는 파일명 생성
                        if st.session_state.generated_slides:
                            first_title = st.session_state.generated_slides[0].get('title', 'presentation')
                            # 특수문자 제거
                            safe_title = re.sub(r'[^\w\s-]', '', first_title).strip()[:20]
                            timestamp = datetime.now().strftime("%m%d")
                            filename = f"{safe_title}_{timestamp}.pptx"
                        else:
                            filename = f"ai_presentation_{datetime.now().strftime('%m%d_%H%M')}.pptx"
                        
                        st.download_button(
                            label="📥 지금 다운로드하기",
                            data=ppt_file,
                            file_name=filename,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True,
                            help="완성된 PPT를 PowerPoint에서 바로 열 수 있습니다!"
                        )
                        
                        st.success("🎉 PPT 파일이 준비되었습니다! 다운로드 버튼을 클릭하세요.")
        
            with col_b:
                # JSON 다운로드 (백업용)
                json_data = json.dumps(st.session_state.generated_slides, ensure_ascii=False, indent=2)
                st.download_button(
                    label="📋 JSON 백업 다운로드",
                    data=json_data,
                    file_name=f"slides_backup_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json",
                    mime="application/json",
                    use_container_width=True,
                    help="나중에 편집할 수 있도록 데이터를 백업하세요"
                )
        else:
            # 빈 상태 표시
            st.markdown(
                """
                <div style='text-align: center; padding: 3rem; background: rgba(102, 126, 234, 0.05); border-radius: 15px;'>
                    <h3 style='color: #667eea;'>🎨 시작하기</h3>
                    <p style='color: #718096; margin-top: 1rem;'>왼쪽에서 PPT 정보를 입력하고<br>생성 버튼을 클릭하세요!</p>
                </div>
                """, 
                unsafe_allow_html=True
            )

# 푸터
st.markdown(
    """
    <div class='footer'>
        <p style='font-size: 1.1rem; font-weight: 500;'>Powered by Kwak Dooil | PPT Generator v1.0</p>
        <p style='font-size: 0.9rem; margin-top: 0.5rem; opacity: 0.8;'>Made with ❤️ for YOU</p>
    </div>
    """, 
    unsafe_allow_html=True
)